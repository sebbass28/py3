from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


BASE = Path(r"c:\Users\Sebas\Desktop\py3")
OUT = BASE / "Presentacion_TFG_DentaLinkLab.pptx"
IMG = BASE / "ilovepdf_images"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0


def add_image_slide(prs, title, image_path, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.3), width=Inches(8.0))
    tx = slide.shapes.add_textbox(Inches(8.7), Inches(1.3), Inches(4.0), Inches(4.8))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {b}"
        p.level = 0


def add_demo_slide(prs):
    bullets = [
        "1) Login por rol y acceso a dashboard",
        "2) Clinica: revision/creacion de pedido",
        "3) Chat por pedido + adjunto de imagen",
        "4) Modo conversacion clinica-laboratorio",
        "5) Laboratorio: cola/kanban de produccion",
        "6) Pacientes: ficha, edicion y anonimizado",
        "7) Finder publico: filtros + comparacion",
        "8) Recuperacion de contraseña funcional",
    ]
    add_bullets_slide(prs, "Demo en vivo (guion 6-8 min)", bullets)


def style_all(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(24 if shape == getattr(slide.shapes, "title", None) else 18)


def main():
    prs = Presentation()

    add_title_slide(
        prs,
        "DentaLinkLab",
        "Proyecto DAW | Juan Sebastian Agudelo Castrillon | IES El Grao | 2024-2025",
    )

    add_bullets_slide(
        prs,
        "Problema detectado",
        [
            "Comunicacion clinica-laboratorio fragmentada en canales no estructurados",
            "Baja trazabilidad del pedido y de las decisiones tecnicas",
            "Dificultad para seguimiento de estados, tiempos y responsabilidades",
        ],
    )

    add_bullets_slide(
        prs,
        "Objetivo del proyecto",
        [
            "Centralizar el flujo completo entre clinica y laboratorio",
            "Aportar trazabilidad, control operativo y mejor experiencia de usuario",
            "Construir una base escalable y desplegable en cloud",
        ],
    )

    add_bullets_slide(
        prs,
        "Alcance del MVP",
        [
            "Auth JWT + roles (clinica/lab)",
            "Pacientes, pedidos, eventos y notificaciones",
            "Chat por pedido y conversacion conjunta",
            "Finder publico de clinicas, metricas y export CSV",
            "Despliegue en AWS con Coolify y Docker",
        ],
    )

    add_image_slide(
        prs,
        "Arquitectura del sistema",
        IMG / "img211.jpg",
        [
            "Frontend React desacoplado del backend",
            "API en Django REST con autenticacion JWT",
            "Persistencia en base de datos relacional",
        ],
    )

    add_image_slide(
        prs,
        "Modelo de datos (ERD)",
        IMG / "img222.jpg",
        [
            "Pedido como entidad central de negocio",
            "Relacion con usuario, paciente, producto y factura",
            "Base para trazabilidad y evolucion futura",
        ],
    )

    add_image_slide(
        prs,
        "Flujo de trabajo CAD/CAM",
        IMG / "img242.jpg",
        [
            "Subida de escaneado y datos del caso",
            "Diseno, validacion, produccion y envio",
            "Trazabilidad de estado y comunicacion",
        ],
    )

    add_image_slide(
        prs,
        "Stack tecnologico",
        IMG / "img257.jpg",
        [
            "Backend: Python + Django",
            "Frontend: React + Tailwind + Three.js",
            "Infraestructura: Docker + AWS",
        ],
    )

    add_image_slide(
        prs,
        "CI/CD y despliegue",
        IMG / "img268.jpg",
        [
            "GitHub como fuente de versionado",
            "Coolify automatiza build, deploy y actualizacion",
            "Produccion en instancia AWS",
        ],
    )

    add_bullets_slide(
        prs,
        "Funcionalidades clave implementadas",
        [
            "Dashboard por rol con filtros y acciones",
            "Chat con adjuntos y mejoras UX",
            "Pacientes: alta, edicion, pedidos y anonimizado",
            "Integraciones MVP con API key y logs",
            "Recuperacion de contraseña funcional",
        ],
    )

    add_demo_slide(prs)

    add_bullets_slide(
        prs,
        "Resultados obtenidos",
        [
            "Mejora de trazabilidad y visibilidad operativa",
            "Centralizacion de informacion del caso",
            "Despliegue real y flujo Scrum demostrado en GitHub",
            "Base lista para evolucion a producto",
        ],
    )

    add_bullets_slide(
        prs,
        "Limitaciones actuales",
        [
            "Pagos no activados en produccion",
            "Notificaciones en polling (sin websocket realtime)",
            "Cobertura E2E ampliable",
        ],
    )

    add_bullets_slide(
        prs,
        "Lineas futuras",
        [
            "Integrar Stripe (sandbox -> produccion)",
            "Migrar notificaciones/chat a tiempo real",
            "Ampliar analitica y observabilidad",
            "Mejorar testing automatizado end-to-end",
        ],
    )

    add_bullets_slide(
        prs,
        "Cierre",
        [
            "DentaLinkLab cumple los objetivos tecnicos y funcionales del modulo",
            "Proyecto desplegado, demostrable y con continuidad real",
            "Gracias. Turno de preguntas",
        ],
    )

    style_all(prs)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()

